"""
Document loaders — extract plain text pages from various file formats.

Each loader returns a list of dicts:
  {text: str, page: int, metadata: dict}

Supported formats:
  .pdf            pymupdf (fitz)
  .docx           python-docx
  .pptx           python-pptx
  .xlsx .xls      openpyxl
  .csv            built-in csv
  .txt .md .rst   plain UTF-8 read
  .msg            extract-msg (Outlook)
  .eml            email (stdlib)
  .png .jpg .jpeg .tiff .bmp  pytesseract OCR (optional — needs Tesseract binary)
"""
from __future__ import annotations
import csv
import email as _email
import io
import logging
import multiprocessing
import threading
from concurrent.futures import ProcessPoolExecutor
from concurrent.futures.process import BrokenProcessPool
from pathlib import Path

_log = logging.getLogger("suni.ingestion.loader")

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".pptx", ".ppt",
    ".xlsx", ".xls", ".csv",
    ".txt", ".md", ".rst", ".log",
    ".msg", ".eml",
    ".png", ".jpg", ".jpeg", ".tiff", ".bmp",
}

from ..system_profile import MAX_DOC_FILE_MB
MAX_FILE_BYTES = MAX_DOC_FILE_MB * 1024 * 1024   # derived from system RAM


# ── Subprocess PDF pool (isolates native fitz crashes from the main process) ──

def _pdf_worker(file_path: str) -> list[dict]:
    """Run fitz in an isolated subprocess — a native crash here stays contained."""
    import fitz  # imported inside worker so spawn doesn't pull the whole module
    pages = []
    with fitz.open(file_path) as doc:
        for i, page in enumerate(doc):
            text = page.get_text("text").strip()
            if text:
                pages.append({"text": text, "page": i + 1, "metadata": {}})
    return pages


_PDF_POOL: ProcessPoolExecutor | None = None
_PDF_POOL_LOCK = threading.Lock()


def _get_pdf_pool() -> ProcessPoolExecutor:
    global _PDF_POOL
    if _PDF_POOL is None:
        with _PDF_POOL_LOCK:
            if _PDF_POOL is None:
                ctx = multiprocessing.get_context("spawn")
                _PDF_POOL = ProcessPoolExecutor(max_workers=2, mp_context=ctx)
    return _PDF_POOL


def load(file_path: str) -> list[dict]:
    """
    Load a document and return a list of page/section dicts.
    Returns [] on unsupported format, oversized file, or unrecoverable error.
    """
    p = Path(file_path)
    try:
        size = p.stat().st_size
        if size > MAX_FILE_BYTES:
            _log.warning("skipping oversized file %s (%.1f MB > 100 MB limit)",
                         file_path, size / 1_048_576)
            return []
    except OSError:
        return []
    ext = p.suffix.lower()
    try:
        if ext == ".pdf":               return _pdf(p)
        if ext in (".docx", ".doc"):    return _docx(p)
        if ext in (".pptx", ".ppt"):    return _pptx(p)
        if ext in (".xlsx", ".xls"):    return _xlsx(p)
        if ext == ".csv":               return _csv(p)
        if ext in (".txt",".md",".rst",".log"): return _text(p)
        if ext == ".msg":               return _msg(p)
        if ext == ".eml":               return _eml(p)
        if ext in (".png",".jpg",".jpeg",".tiff",".bmp"): return _ocr(p)
    except Exception as e:
        _log.warning("load failed %s: %s", file_path, e)
    return []


# ── Loaders ───────────────────────────────────────────────────────────────────

def _pdf(p: Path) -> list[dict]:
    global _PDF_POOL
    for attempt in range(2):
        try:
            pool = _get_pdf_pool()
            return pool.submit(_pdf_worker, str(p)).result(timeout=120)
        except BrokenProcessPool:
            # Worker process crashed (e.g. native fitz fault) — reset pool and skip file
            _log.warning("pdf worker crashed on %s — rebuilding pool", p)
            with _PDF_POOL_LOCK:
                _PDF_POOL = None
            if attempt == 1:
                return []
        except Exception as e:
            _log.warning("pdf load failed %s: %s", p, e)
            return []
    return []


def _docx(p: Path) -> list[dict]:
    from docx import Document
    doc = Document(str(p))
    # Group paragraphs into ~sections (split on Heading styles)
    sections, buf, section_num = [], [], 1
    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        if para.style.name.startswith("Heading"):
            if buf:
                sections.append({"text": "\n".join(buf), "page": section_num, "metadata": {}})
                section_num += 1
                buf = []
        buf.append(t)
    if buf:
        sections.append({"text": "\n".join(buf), "page": section_num, "metadata": {}})
    return sections


def _pptx(p: Path) -> list[dict]:
    from pptx import Presentation
    slides = []
    prs = Presentation(str(p))
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides.append({"text": "\n".join(parts), "page": i + 1, "metadata": {}})
    return slides


def _xlsx(p: Path) -> list[dict]:
    from openpyxl import load_workbook
    wb = load_workbook(str(p), read_only=True, data_only=True)
    sheets = []
    for sheet_num, ws in enumerate(wb.worksheets, 1):
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append("\t".join(cells))
        if rows:
            sheets.append({
                "text": "\n".join(rows),
                "page": sheet_num,
                "metadata": {"sheet": ws.title},
            })
    wb.close()
    return sheets


def _csv(p: Path) -> list[dict]:
    try:
        text = p.read_text(encoding="utf-8", errors="replace")
    except Exception:
        return []
    rows = []
    for row in csv.reader(io.StringIO(text)):
        if any(c.strip() for c in row):
            rows.append("\t".join(row))
    if not rows:
        return []
    # Chunk every 200 rows into one page
    chunk_size = 200
    pages = []
    for i in range(0, len(rows), chunk_size):
        pages.append({"text": "\n".join(rows[i:i+chunk_size]), "page": i//chunk_size+1, "metadata": {}})
    return pages


def _text(p: Path) -> list[dict]:
    try:
        text = p.read_text(encoding="utf-8", errors="replace").strip()
    except Exception:
        return []
    if not text:
        return []
    return [{"text": text, "page": 1, "metadata": {}}]


def _msg(p: Path) -> list[dict]:
    try:
        import extract_msg
        msg = extract_msg.Message(str(p))
        parts = []
        if msg.subject:   parts.append(f"Subject: {msg.subject}")
        if msg.sender:    parts.append(f"From: {msg.sender}")
        if msg.date:      parts.append(f"Date: {msg.date}")
        if msg.body:      parts.append(msg.body.strip())
        msg.close()
        return [{"text": "\n".join(parts), "page": 1, "metadata": {}}] if parts else []
    except Exception as e:
        _log.warning("msg load failed %s: %s", p, e)
        return []


def _eml(p: Path) -> list[dict]:
    try:
        raw = p.read_bytes()
        msg = _email.message_from_bytes(raw)
        parts = []
        if msg.get("subject"):  parts.append(f"Subject: {msg.get('subject')}")
        if msg.get("from"):     parts.append(f"From: {msg.get('from')}")
        if msg.get("date"):     parts.append(f"Date: {msg.get('date')}")
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    charset = part.get_content_charset() or "utf-8"
                    parts.append(payload.decode(charset, errors="replace").strip())
        return [{"text": "\n".join(parts), "page": 1, "metadata": {}}] if parts else []
    except Exception:
        return []


def _ocr(p: Path) -> list[dict]:
    try:
        import pytesseract
        from PIL import Image
        text = pytesseract.image_to_string(Image.open(str(p))).strip()
        return [{"text": text, "page": 1, "metadata": {}}] if text else []
    except ImportError:
        _log.debug("pytesseract/Pillow not installed — skipping image %s", p)
        return []
    except Exception as e:
        _log.warning("OCR failed %s: %s", p, e)
        return []
